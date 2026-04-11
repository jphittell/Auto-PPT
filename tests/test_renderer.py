from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor

from pptx_gen.layout.schemas import ResolvedDeckLayout, StyleTokens
from pptx_gen.renderer.pptx_exporter import export_pptx
from pptx_gen.renderer.qa import QAStatus, validate_export, validate_layout
from pptx_gen.renderer.slide_ops import add_image, add_shape, add_text, set_background_color


def test_slide_ops_create_text_image_and_shape(tmp_path: Path, style_tokens_payload: dict, tiny_png_bytes: bytes) -> None:
    image_path = tmp_path / "tiny.png"
    image_path.write_bytes(tiny_png_bytes)
    tokens = StyleTokens(**style_tokens_payload)

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    set_background_color(slide, tokens.colors.bg)
    add_text(slide, "Hello world", x=0.5, y=0.5, w=4.0, h=1.0, style_tokens=tokens, style_ref="body")
    add_image(slide, image_path, x=0.5, y=1.8, w=1.0, h=1.0)
    add_shape(slide, x=2.0, y=1.8, w=2.0, h=1.0, style_tokens=tokens, style_ref="takeaway", text="Key point")

    assert len(slide.shapes) == 3
    assert slide.shapes[0].has_text_frame
    assert slide.background.fill.fore_color.rgb == RGBColor.from_string(tokens.colors.bg.removeprefix("#"))


def test_add_image_requires_existing_local_asset(style_tokens_payload: dict) -> None:
    tokens = StyleTokens(**style_tokens_payload)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(FileNotFoundError):
        add_image(slide, "C:/missing/image.png", x=0.5, y=0.5, w=1.0, h=1.0)

    add_shape(slide, x=1.5, y=0.5, w=1.0, h=0.5, style_tokens=tokens, style_ref="accent_bar")


def test_export_pptx_and_export_qa(tmp_path: Path, style_tokens_payload: dict, tiny_png_bytes: bytes) -> None:
    image_path = tmp_path / "chart.png"
    image_path.write_bytes(tiny_png_bytes)
    tokens = StyleTokens(**style_tokens_payload)
    layout = ResolvedDeckLayout(
        deck_id="deck-1",
        slides=[
            {
                "slide_id": "s1",
                "elements": [
                    {
                        "element_id": "s1:headline",
                        "kind": "textbox",
                        "x": 0.75,
                        "y": 0.75,
                        "w": 6.0,
                        "h": 0.75,
                        "z": 0,
                        "data_ref": "slide:s1:headline",
                        "style_ref": "headline",
                        "payload": {"content": "Quarterly Review"},
                    },
                    {
                        "element_id": "s1:chart",
                        "kind": "image",
                        "x": 0.75,
                        "y": 1.75,
                        "w": 3.0,
                        "h": 3.0,
                        "z": 0,
                        "data_ref": "block:b1",
                        "style_ref": "image",
                        "payload": {"content": {"path": str(image_path), "width_px": 1, "height_px": 1}},
                    },
                ],
            }
        ],
    )

    output_path = export_pptx(layout=layout, style_tokens=tokens, output_path=tmp_path / "deck.pptx")
    report = validate_export(output_path, layout=layout, style_tokens=tokens)

    assert output_path.exists()
    assert report.passed is True
    assert report.executability.status is QAStatus.PASS


def test_export_pptx_removes_existing_template_slides(
    tmp_path: Path,
    style_tokens_payload: dict,
    make_pptx_file,
) -> None:
    tokens = StyleTokens(**style_tokens_payload)
    template_path = make_pptx_file()
    layout = ResolvedDeckLayout(
        deck_id="deck-template-cleanup",
        slides=[
            {
                "slide_id": "s1",
                "elements": [
                    {
                        "element_id": "s1:headline",
                        "kind": "textbox",
                        "x": 0.75,
                        "y": 0.75,
                        "w": 6.0,
                        "h": 0.75,
                        "z": 0,
                        "data_ref": "slide:s1:headline",
                        "style_ref": "headline",
                        "payload": {"content": "Generated Title"},
                    }
                ],
            }
        ],
    )

    output_path = export_pptx(
        layout=layout,
        style_tokens=tokens,
        output_path=tmp_path / "deck-from-template.pptx",
        template_path=template_path,
    )
    exported = Presentation(str(output_path))
    first_slide = exported.slides[0]
    text_content = "\n".join(
        shape.text for shape in first_slide.shapes if getattr(shape, "has_text_frame", False)
    )

    assert len(exported.slides) == 1
    assert "Generated Title" in text_content


def test_layout_qa_reports_overlap_bounds_contrast_and_missing_assets(style_tokens_payload: dict) -> None:
    low_contrast_tokens = StyleTokens(
        **{
            **style_tokens_payload,
            "colors": {
                **style_tokens_payload["colors"],
                "muted": "#A0A0A0",
            },
        }
    )
    layout = ResolvedDeckLayout(
        deck_id="deck-qa",
        slides=[
            {
                "slide_id": "s1",
                "elements": [
                    {
                        "element_id": "s1:text1",
                        "kind": "textbox",
                        "x": 0.75,
                        "y": 0.75,
                        "w": 2.0,
                        "h": 0.5,
                        "z": 0,
                        "data_ref": "b1",
                        "style_ref": "subtitle",
                        "payload": {"content": {"text": " ".join(["word"] * 50)}},
                    },
                    {
                        "element_id": "s1:text2",
                        "kind": "textbox",
                        "x": 1.0,
                        "y": 0.9,
                        "w": 2.0,
                        "h": 0.5,
                        "z": 0,
                        "data_ref": "b2",
                        "style_ref": "body",
                        "payload": {"content": {"text": "overlap block"}},
                    },
                    {
                        "element_id": "s1:image",
                        "kind": "image",
                        "x": 12.8,
                        "y": 6.9,
                        "w": 1.0,
                        "h": 1.0,
                        "z": 0,
                        "data_ref": "b3",
                        "style_ref": "image",
                        "payload": {"content": {"path": "C:/missing/asset.png"}},
                    },
                ],
            }
        ],
    )

    report = validate_layout(layout, style_tokens=low_contrast_tokens)

    assert report.passed is False
    assert any(item.check == "layout.overlap" and item.status is QAStatus.FAIL for item in report.items)
    assert any(item.check == "layout.bounds" and item.status is QAStatus.FAIL for item in report.items)
    assert any(item.check == "color.contrast" for item in report.items)
    assert any(item.check == "image.local_asset" and item.status is QAStatus.FAIL for item in report.items)
    assert any(item.check == "text.overflow" and item.status is QAStatus.FAIL for item in report.items)
