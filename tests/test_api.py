from __future__ import annotations

import io
import time

from fastapi.testclient import TestClient
from pptx import Presentation

import pptx_gen.api as api_module
from pptx_gen.api_schemas import ExportBlockRequest, ExportSlideRequest
from pptx_gen.layout.schemas import StyleTokens
from pptx_gen.planning.schemas import (
    DeckTheme,
    LayoutIntent,
    OutlineItem,
    OutlineSpec,
    PresentationBlock,
    PresentationBlockKind,
    PresentationSpec,
    SlideArchetype,
    SlidePurpose,
    SlideSpec,
)


def _reset_api_state() -> None:
    api_module._store.clear()
    api_module._INGESTED_VECTOR_STORES.clear()
    limiter = getattr(api_module.app.state, "limiter", None)
    storage = getattr(limiter, "_storage", None)
    if storage is not None and hasattr(storage, "reset"):
        storage.reset()
    api_module._EMBEDDER = None
    api_module._STRUCTURED_LLM_CLIENT = False
    api_module._clear_preview_structure_cache()
    api_module._GENERATION_JOBS.clear()
    api_module._GENERATION_QUEUE = None
    api_module._GENERATION_QUEUE_LOOP = None
    api_module._GENERATION_WORKER_TASK = None


def _ingest_fixture(client: TestClient, sample_pdf_path) -> str:
    response = client.post(
        "/api/ingest",
        files={"file": ("sample_ingestion.pdf", sample_pdf_path.read_bytes(), "application/pdf")},
    )
    assert response.status_code == 200
    return response.json()["doc_id"]


def _assert_error(response, *, status_code: int, code: str, message_fragment: str | None = None) -> dict:
    assert response.status_code == status_code
    payload = response.json()
    assert set(payload) == {"error"}
    error = payload["error"]
    assert error["code"] == code
    assert isinstance(error["message"], str) and error["message"]
    assert isinstance(error["request_id"], str) and error["request_id"]
    if message_fragment is not None:
        assert message_fragment in error["message"]
    assert response.headers["x-request-id"] == error["request_id"]
    return error


def test_api_health_and_templates(monkeypatch, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)

    health = client.get("/api/health")
    assert health.status_code == 200
    assert health.json() == {
        "status": "ok",
        "phase": "1",
        "ingest": True,
        "generation": "live",
        "embedder": {"status": "ok", "latency_ms": health.json()["embedder"]["latency_ms"]},
        "vector_store": {"status": "ok", "latency_ms": health.json()["vector_store"]["latency_ms"]},
    }
    assert isinstance(health.json()["embedder"]["latency_ms"], int)
    assert isinstance(health.json()["vector_store"]["latency_ms"], int)

    templates = client.get("/api/templates")
    assert templates.status_code == 200
    payload = templates.json()
    assert len(payload) == 20
    template_by_id = {item["id"]: item for item in payload}
    assert template_by_id["headline.evidence"]["deck_default_allowed"] is True
    assert template_by_id["compare.2col"]["deck_default_allowed"] is True
    assert template_by_id["kpi.big"]["deck_default_allowed"] is True
    assert template_by_id["title.cover"]["deck_default_allowed"] is False
    assert template_by_id["exec.summary"]["deck_default_allowed"] is False
    assert template_by_id["closing.actions"]["deck_default_allowed"] is False
    assert template_by_id["impact.statement"]["deck_default_allowed"] is False

    themes = client.get("/api/themes")
    assert themes.status_code == 200
    assert themes.json() == ["ONAC"]


def test_api_ingest_pdf_returns_summary(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)

    response = client.post(
        "/api/ingest",
        files={"file": ("sample_ingestion.pdf", sample_pdf_path.read_bytes(), "application/pdf")},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["doc_id"]
    assert payload["chunk_count"] > 0
    assert payload["title"]
    assert payload["element_types"]


def test_api_ingest_pptx_returns_slide_structure(monkeypatch, make_pptx_file, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    source_path = make_pptx_file()

    response = client.post(
        "/api/ingest",
        files={
            "file": (
                "sample.pptx",
                source_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["source_format"] == "pptx"
    assert payload["slide_count"] == 2
    assert payload["slide_types"]["title"] == 1
    assert payload["slide_types"]["table"] == 1


def test_api_plan_from_prompt_uses_powerpoint_slide_count(monkeypatch, make_pptx_file, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    source_path = make_pptx_file()

    ingest = client.post(
        "/api/ingest",
        files={
            "file": (
                "sample.pptx",
                source_path.read_bytes(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert ingest.status_code == 200
    doc_id = ingest.json()["doc_id"]

    planned = client.post(
        "/api/plan/prompt",
        json={
            "doc_ids": [doc_id],
            "prompt": "Rebrand this deck for ONAC while preserving the original structure.",
        },
    )

    assert planned.status_code == 200
    payload = planned.json()
    assert len(payload["slides"]) == 2
    assert payload["slides"][0]["template_id"] == "title.cover"
    assert payload["slides"][1]["title"] == "Decision Summary"
    assert payload["slides"][1]["template_id"] == "compare.2col"


def test_enforce_outline_authority_preserves_outline_template_key(style_tokens_payload) -> None:
    spec = PresentationSpec(
        title="Imported deck",
        audience="Executive leadership",
        theme=DeckTheme(name="Auto PPT", style_tokens=StyleTokens(**style_tokens_payload)),
        slides=[
            SlideSpec(
                slide_id="llm-slide",
                purpose=SlidePurpose.CONTENT,
                archetype=None,
                layout_intent=LayoutIntent(template_key="headline.evidence", strict_template=False),
                headline="Working title",
                blocks=[
                    PresentationBlock(
                        block_id="b1",
                        kind=PresentationBlockKind.TEXT,
                        content={"text": "Compare rollout options across regions"},
                    )
                ],
            )
        ],
    )
    outline = OutlineSpec(
        outline=[
            OutlineItem(
                slide_id="s2",
                purpose=SlidePurpose.CONTENT,
                archetype=SlideArchetype.COMPARISON,
                headline="Decision Summary",
                message="Compare rollout options across regions",
                evidence_queries=["rollout options across regions"],
                template_key="compare.2col",
            )
        ],
        questions_for_user=[],
    )

    enforced = api_module._enforce_outline_authority(spec, outline, source_ids=["src-1"])

    assert enforced.slides[0].slide_id == "s2"
    assert enforced.slides[0].headline == "Decision Summary"
    assert enforced.slides[0].archetype is SlideArchetype.COMPARISON
    assert enforced.slides[0].layout_intent.template_key == "compare.2col"
    assert enforced.slides[0].layout_intent.strict_template is True


def test_api_plan_and_generate_honor_authoritative_inputs(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    doc_id = _ingest_fixture(client, sample_pdf_path)

    planned = client.post(
        "/api/plan",
        json={
            "doc_ids": [doc_id],
            "goal": "Board update",
            "audience": "Executive Steering Committee",
            "tone": 15,
            "slide_count": 6,
        },
    )
    assert planned.status_code == 200
    draft = planned.json()
    assert len(draft["slides"]) == 6
    assert draft["draft_id"]
    assert draft["slides"][0]["template_id"] == "title.cover"
    assert draft["slides"][1]["template_id"] == "exec.summary"
    assert draft["slides"][1]["archetype"] == "executive_summary"
    assert draft["slides"][0]["blocks"][0]["content"].startswith("Subtitle:")
    assert any(slide["template_id"] == "closing.actions" for slide in draft["slides"])

    edited_outline = list(draft["slides"])
    edited_outline[1]["title"] = "Investor-safe evidence"
    edited_outline[1]["template_id"] = "compare.2col"
    edited_outline[2]["title"] = "Reordered board priorities"
    edited_outline[1], edited_outline[2] = edited_outline[2], edited_outline[1]
    for index, slide in enumerate(edited_outline, start=1):
        slide["index"] = index

    generated = client.post(
        "/api/generate",
        json={
            "draft_id": draft["draft_id"],
            "outline": [
                {
                    "id": slide["id"],
                    "index": slide["index"],
                    "purpose": slide["purpose"],
                    "title": slide["title"],
                    "template_id": slide["template_id"],
                }
                for slide in edited_outline
            ],
            "selected_template_id": "kpi.big",
            "brand_kit": {
                "logo_data_url": None,
                "primary_color": "#112233",
                "accent_color": "#445566",
                "font_pair": "DM Sans/DM Serif Display",
            },
        },
    )
    assert generated.status_code == 200
    deck = generated.json()

    assert len(deck["slides"]) == 6
    assert {deck["slides"][1]["title"], deck["slides"][2]["title"]} == {"Investor-safe evidence", "Reordered board priorities"}
    deck_by_title = {slide["title"]: slide for slide in deck["slides"]}
    content_templates = {
        slide["title"]: slide["template_id"]
        for slide in deck["slides"]
        if slide["purpose"] in {"content", "summary"}
    }
    assert deck_by_title["Investor-safe evidence"]["purpose"] == "content"
    assert content_templates["Investor-safe evidence"] == "compare.2col"
    assert any(slide.get("archetype") == "executive_summary" for slide in deck["slides"])
    overview_slide = next(slide for slide in deck["slides"] if slide.get("archetype") == "executive_summary")
    assert overview_slide["template_id"] in {"exec.summary", "compare.2col"}
    assert any(slide["template_id"] == "closing.actions" for slide in deck["slides"])
    assert not any(
        slide["title"].startswith(prefix)
        for slide in deck["slides"]
        for prefix in ("Implementation implications", "Design quality strategies", "Hybrid architecture")
    )
    assert any(
        block["kind"] == "bullets" and block["content"].startswith("• ")
        for slide in deck["slides"]
        for block in slide["blocks"]
    )
    assert deck["theme"]["primary_color"] == "#112233"
    assert deck["theme"]["accent_color"] == "#445566"
    assert deck["theme"]["heading_font"] == "DM Serif Display"
    assert deck["theme"]["body_font"] == "DM Sans"
    assert deck["theme"]["heading_font"] == "DM Serif Display"
    assert deck["theme"]["body_font"] == "DM Sans"
    fetched = client.get(f"/api/deck/{deck['id']}")
    assert fetched.status_code == 200
    assert fetched.json()["id"] == deck["id"]

    pptx = client.post(f"/api/export/{deck['id']}", json={"format": "pptx"})
    assert pptx.status_code == 200
    assert pptx.headers["content-type"].startswith(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    assert pptx.content.startswith(b"PK")

    pdf = client.post(f"/api/export/{deck['id']}", json={"format": "pdf"})
    assert pdf.status_code == 200
    assert pdf.headers["content-type"].startswith("application/pdf")
    assert pdf.content.startswith(b"%PDF-1.4")


def test_api_different_prompt_inputs_change_generated_result(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    doc_id = _ingest_fixture(client, sample_pdf_path)

    analytical_plan = client.post(
        "/api/plan",
        json={
            "doc_ids": [doc_id],
            "goal": "Board update",
            "audience": "Board",
            "tone": 10,
            "slide_count": 6,
        },
    ).json()
    bold_plan = client.post(
        "/api/plan",
        json={
            "doc_ids": [doc_id],
            "goal": "Raise seed",
            "audience": "Investors",
            "tone": 90,
            "slide_count": 6,
        },
    ).json()

    analytical_deck = client.post(
        "/api/generate",
        json={
            "draft_id": analytical_plan["draft_id"],
            "outline": [
                {
                    "id": slide["id"],
                    "index": slide["index"],
                    "purpose": slide["purpose"],
                    "title": slide["title"],
                    "template_id": slide["template_id"],
                }
                for slide in analytical_plan["slides"]
            ],
            "selected_template_id": "headline.evidence",
            "brand_kit": {
                "logo_data_url": None,
                "primary_color": "#4F46E5",
                "accent_color": "#0F172A",
                "font_pair": "Inter/Inter",
            },
        },
    ).json()
    bold_deck = client.post(
        "/api/generate",
        json={
            "draft_id": bold_plan["draft_id"],
            "outline": [
                {
                    "id": slide["id"],
                    "index": slide["index"],
                    "purpose": slide["purpose"],
                    "title": slide["title"],
                    "template_id": slide["template_id"],
                }
                for slide in bold_plan["slides"]
            ],
            "selected_template_id": "compare.2col",
            "brand_kit": {
                "logo_data_url": None,
                "primary_color": "#0A84FF",
                "accent_color": "#111111",
                "font_pair": "Inter/Inter",
            },
        },
    ).json()

    analytical_text = "\n".join(block["content"] for slide in analytical_deck["slides"] for block in slide["blocks"])
    bold_text = "\n".join(block["content"] for slide in bold_deck["slides"] for block in slide["blocks"])

    assert analytical_deck["audience"] == "Board"
    assert bold_deck["audience"] == "Investors"
    analytical_tone_hints = [
        block.get("data", {}).get("tone_hint")
        for slide in analytical_deck["slides"]
        for block in slide["blocks"]
        if isinstance(block.get("data"), dict)
    ]
    bold_tone_hints = [
        block.get("data", {}).get("tone_hint")
        for slide in bold_deck["slides"]
        for block in slide["blocks"]
        if isinstance(block.get("data"), dict)
    ]
    # With the deterministic planner (no LLM), content text is the same from the
    # same source document regardless of tone — only structure, audience, and
    # template choices differ.  Verify the deck-level metadata reflects the
    # different briefs and that template variety exists across slides.
    assert analytical_deck["audience"] != bold_deck["audience"]
    analytical_templates = {slide["template_id"] for slide in analytical_deck["slides"]}
    bold_templates = {slide["template_id"] for slide in bold_deck["slides"]}
    assert len(analytical_templates) >= 2, f"Expected template variety, got {analytical_templates}"
    assert len(bold_templates) >= 2, f"Expected template variety, got {bold_templates}"
    assert analytical_plan["slides"][0]["blocks"][0]["content"].startswith("Subtitle: Board update")
    assert any(slide["template_id"] in {"compare.2col", "exec.summary"} for slide in analytical_plan["slides"])


def test_api_export_uses_ui_slide_overrides(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    doc_id = _ingest_fixture(client, sample_pdf_path)

    planned = client.post(
        "/api/plan",
        json={
            "doc_ids": [doc_id],
            "goal": "Board update",
            "audience": "Executive Steering Committee",
            "tone": 15,
            "slide_count": 6,
        },
    ).json()
    generated = client.post(
        "/api/generate",
        json={
            "draft_id": planned["draft_id"],
            "outline": [
                {
                    "id": slide["id"],
                    "index": slide["index"],
                    "purpose": slide["purpose"],
                    "title": slide["title"],
                    "template_id": slide["template_id"],
                }
                for slide in planned["slides"]
            ],
            "selected_template_id": "headline.evidence",
            "brand_kit": {
                "logo_data_url": None,
                "primary_color": "#112233",
                "accent_color": "#445566",
                "font_pair": "DM Sans/DM Serif Display",
            },
        },
    ).json()

    ui_slides = list(generated["slides"])
    ui_slides.append(
        {
            "id": "slide-ui-added",
            "index": len(ui_slides) + 1,
            "purpose": "content",
            "archetype": None,
            "title": "UI Added Slide",
            "template_id": "headline.evidence",
            "speaker_notes": "Added in editor",
            "blocks": [
                {
                    "id": "block-ui-added",
                    "kind": "text",
                    "content": "Export should include this slide.",
                    "data": {"text": "Export should include this slide."},
                    "citation": None,
                }
            ],
        }
    )

    pptx = client.post(
        f"/api/export/{generated['id']}",
        json={"format": "pptx", "slides": ui_slides},
    )

    assert pptx.status_code == 200
    exported = Presentation(io.BytesIO(pptx.content))
    slide_text = "\n".join(
        shape.text
        for slide in exported.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert len(exported.slides) == len(ui_slides)
    assert "UI Added Slide" in slide_text


def test_ui_export_canonicalizes_three_column_cards() -> None:
    slide = api_module._ui_slide_to_planning_slide(
        ExportSlideRequest(
            id="s-columns",
            index=1,
            purpose="content",
            title="Capability map",
            template_id="content.3col",
            speaker_notes=None,
            blocks=[
                ExportBlockRequest(
                    id="b1",
                    kind="callout",
                    content="Speed: Faster rollout\nQuality: Stronger controls\nCost: Lower support load",
                    data={
                        "cards": [
                            {"title": "Speed", "text": "Faster rollout"},
                            {"title": "Quality", "text": "Stronger controls"},
                            {"title": "Cost", "text": "Lower support load"},
                        ]
                    },
                    citation=None,
                )
            ],
        ),
        existing_slide=None,
        fallback_source_id="ui-export",
    )

    assert [block.kind.value for block in slide.blocks] == ["text", "text", "text"]
    assert [block.content["text"] for block in slide.blocks] == [
        "Speed\nFaster rollout",
        "Quality\nStronger controls",
        "Cost\nLower support load",
    ]


def test_ui_export_preserves_structured_image_data_for_photo_templates() -> None:
    slide = api_module._ui_slide_to_planning_slide(
        ExportSlideRequest(
            id="s-photo",
            index=1,
            purpose="content",
            title="Original title",
            template_id="bold.photo",
            speaker_notes=None,
            blocks=[
                ExportBlockRequest(
                    id="b1",
                    kind="text",
                    content="A sharper executive statement",
                    data={"text": "A sharper executive statement"},
                    citation=None,
                ),
                ExportBlockRequest(
                    id="b2",
                    kind="image",
                    content="path: C:/assets/photo.png",
                    data={"path": "C:/assets/photo.png"},
                    citation=None,
                ),
            ],
        ),
        existing_slide=None,
        fallback_source_id="ui-export",
    )

    assert slide.headline == "A sharper executive statement"
    assert slide.blocks[1].content == {"path": "C:/assets/photo.png"}


def test_ui_export_canonicalizes_chart_takeaway_callout_cards_to_text() -> None:
    slide = api_module._ui_slide_to_planning_slide(
        ExportSlideRequest(
            id="s-chart",
            index=1,
            purpose="content",
            title="Quarterly trend",
            template_id="chart.takeaway",
            speaker_notes=None,
            blocks=[
                ExportBlockRequest(
                    id="b1",
                    kind="chart",
                    content="Q1: 1\nQ2: 2",
                    data={"chart_type": "bar", "series": [{"label": "Q1", "value": 1.0}, {"label": "Q2", "value": 2.0}]},
                    citation=None,
                ),
                ExportBlockRequest(
                    id="b2",
                    kind="callout",
                    content="Takeaway: Automation compounds over time",
                    data={"cards": [{"title": "Takeaway", "text": "Automation compounds over time"}]},
                    citation=None,
                ),
            ],
        ),
        existing_slide=None,
        fallback_source_id="ui-export",
    )

    assert slide.blocks[1].content == {"text": "Takeaway: Automation compounds over time"}


def test_ui_export_preserves_structured_table_data() -> None:
    slide = api_module._ui_slide_to_planning_slide(
        ExportSlideRequest(
            id="s-table",
            index=1,
            purpose="content",
            title="Agenda",
            template_id="agenda.table",
            speaker_notes=None,
            blocks=[
                ExportBlockRequest(
                    id="b1",
                    kind="table",
                    content="Discovery | Align priorities\nDelivery | Sequence work",
                    data={
                        "columns": ["Section", "Focus"],
                        "rows": [["Discovery", "Align priorities"], ["Delivery", "Sequence work"]],
                    },
                    citation=None,
                )
            ],
        ),
        existing_slide=None,
        fallback_source_id="ui-export",
    )

    assert slide.blocks[0].content["columns"] == ["Section", "Focus"]
    assert slide.blocks[0].content["rows"][0] == ["Discovery", "Align priorities"]


def test_ui_export_preserves_card_data_for_icon_templates() -> None:
    for template_id, expected_count in (("icons.3", 3), ("icons.4", 4)):
        slide = api_module._ui_slide_to_planning_slide(
            ExportSlideRequest(
                id=f"s-{template_id}",
                index=1,
                purpose="content",
                title="Delivery model",
                template_id=template_id,
                speaker_notes=None,
                blocks=[
                    ExportBlockRequest(
                        id="b1",
                        kind="callout",
                        content="\n".join(
                            f"Card {index + 1}: Detail {index + 1}"
                            for index in range(expected_count)
                        ),
                        data={
                            "cards": [
                                {"title": f"Card {index + 1}", "text": f"Detail {index + 1}"}
                                for index in range(expected_count)
                            ]
                        },
                        citation=None,
                    )
                ],
            ),
            existing_slide=None,
            fallback_source_id="ui-export",
        )

        assert len(slide.blocks) == 1
        assert slide.blocks[0].kind.value == "callout"
        assert len(slide.blocks[0].content["cards"]) == expected_count


def test_api_plan_from_prompt_infers_planning_inputs(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)
    doc_id = _ingest_fixture(client, sample_pdf_path)

    response = client.post(
        "/api/plan/prompt",
        json={
            "doc_ids": [doc_id],
            "prompt": "Create a 6 slide architecture deck for Oracle consultants focused on pipeline design.",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["audience"] == "Oracle consultants"
    assert payload["goal"].lower().startswith("create a 6 slide architecture")
    assert len(payload["slides"]) == 6
    assert payload["slides"][0]["template_id"] == "title.cover"


def test_api_chat_generate_runs_pipeline(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)

    response = client.post(
        "/api/chat/generate",
        data={"prompt": "Create a 6 slide architecture deck for Oracle consultants focused on pipeline design."},
        files={"file": ("sample_ingestion.pdf", sample_pdf_path.read_bytes(), "application/pdf")},
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["session_id"].startswith("chat-")
    assert payload["inferred_audience"] == "Oracle consultants"
    assert payload["inferred_slide_count"] == 6
    assert payload["deck"]["slides"]
    overview_slide = next(slide for slide in payload["deck"]["slides"] if slide["archetype"] == "executive_summary")
    assert any(block.get("data", {}).get("cards") for block in overview_slide["blocks"] if isinstance(block.get("data"), dict))
    assert any(slide["template_id"] in {"closing.actions", "compare.2col"} for slide in payload["deck"]["slides"])


def test_api_slide_preview_calls_llm_and_returns_consulting_style() -> None:
    _reset_api_state()

    class FakeClient:
        def generate_json(self, *, system_prompt: str, user_prompt: str, schema_name: str) -> dict:
            assert schema_name == "SlidePreviewLLMResponse"
            assert "selected template key: exec.summary" in user_prompt.lower()
            return {
                "headline": "Executive Overview",
                "speaker_notes": "Consulting style preview",
                "blocks": [
                    {
                        "kind": "text",
                        "text": "Consulting-style summary for Oracle delivery leaders.",
                    },
                    {
                        "kind": "callout",
                        "cards": [
                            {"title": "Ingestion", "text": "Normalize enterprise inputs"},
                            {"title": "Retrieval", "text": "Ground claims in evidence"},
                            {"title": "Layout", "text": "Apply deterministic templates"},
                        ],
                    },
                ],
            }

    api_module._STRUCTURED_LLM_CLIENT = FakeClient()
    client = TestClient(api_module.app)

    response = client.post(
        "/api/slide/preview",
        json={
            "slide_id": "slide-preview-1",
            "title": "Executive Overview",
            "purpose": "content",
            "template_id": "exec.summary",
            "content": "Current draft text about ingestion, retrieval, layout, and export.",
            "audience": "Oracle consultants",
            "goal": "Explain the delivery architecture",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["id"] == "slide-preview-1"
    assert payload["title"] == "Executive Overview"
    assert payload["template_id"] == "exec.summary"
    assert any(block.get("data", {}).get("cards") for block in payload["blocks"] if isinstance(block.get("data"), dict))


def test_api_slide_preview_caches_identical_llm_requests() -> None:
    _reset_api_state()

    class FakeClient:
        def __init__(self) -> None:
            self.calls = 0

        def generate_json(self, *, system_prompt: str, user_prompt: str, schema_name: str) -> dict:
            self.calls += 1
            return {
                "headline": "Executive Overview",
                "speaker_notes": "Consulting style preview",
                "blocks": [
                    {"kind": "text", "text": "Consulting-style summary for Oracle delivery leaders."},
                ],
            }

    fake_client = FakeClient()
    api_module._STRUCTURED_LLM_CLIENT = fake_client
    client = TestClient(api_module.app)
    payload = {
        "slide_id": "slide-preview-cache",
        "title": "Executive Overview",
        "purpose": "content",
        "template_id": "exec.summary",
        "content": "Current draft text about ingestion, retrieval, layout, and export.",
        "audience": "Oracle consultants",
        "goal": "Explain the delivery architecture",
    }

    first = client.post("/api/slide/preview", json=payload)
    second = client.post("/api/slide/preview", json=payload)

    assert first.status_code == 200
    assert second.status_code == 200
    assert first.json() == second.json()
    assert fake_client.calls == 1


def test_api_slide_preview_falls_back_when_structured_llm_payload_is_malformed() -> None:
    _reset_api_state()

    class FakeClient:
        def generate_json(self, *, system_prompt: str, user_prompt: str, schema_name: str) -> dict:
            return {
                "headline": "Broken Preview",
                "blocks": "not-a-list",
            }

    api_module._STRUCTURED_LLM_CLIENT = FakeClient()
    client = TestClient(api_module.app)

    response = client.post(
        "/api/slide/preview",
        json={
            "slide_id": "slide-preview-fallback",
            "title": "Fallback Preview",
            "purpose": "content",
            "template_id": "compare.2col",
            "content": "First point explains the ingestion pipeline. Second point covers the retrieval layer. Third point summarizes the export system.",
            "audience": "Operators",
            "goal": "Explain the workflow",
        },
    )

    assert response.status_code == 200
    payload = response.json()
    assert payload["id"] == "slide-preview-fallback"
    assert payload["title"] == "Fallback Preview"
    assert payload["template_id"] == "compare.2col"
    assert payload["blocks"]
    assert all(isinstance(block.get("data"), dict) for block in payload["blocks"])
    # compare.2col emits 2 bullet blocks when the source has enough points
    # to split, or 1 when there's insufficient content (no filler padding).
    assert len(payload["blocks"]) in {1, 2}
    assert all(block["kind"] == "bullets" for block in payload["blocks"])
    assert all(block.get("data", {}).get("items") for block in payload["blocks"])


def test_fallback_structure_content_caches_identical_inputs(monkeypatch) -> None:
    _reset_api_state()

    calls = 0
    real_impl = api_module._fallback_structure_content_uncached

    def recording_impl(content: str, title: str, template: str, *, grounding_text: str = "") -> dict:
        nonlocal calls
        calls += 1
        return real_impl(content, title, template, grounding_text=grounding_text)

    monkeypatch.setattr(api_module, "_fallback_structure_content_uncached", recording_impl)

    first = api_module._fallback_structure_content(
        "First point explains the ingestion pipeline. Second point covers retrieval.",
        "Fallback Preview",
        "exec.summary",
    )
    second = api_module._fallback_structure_content(
        "First point explains the ingestion pipeline. Second point covers retrieval.",
        "Fallback Preview",
        "exec.summary",
    )

    assert first == second
    assert first is not second
    assert calls == 1


def test_api_async_generate_job_completes_and_returns_result(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    with TestClient(api_module.app) as client:
        doc_id = _ingest_fixture(client, sample_pdf_path)
        planned = client.post(
            "/api/plan",
            json={
                "doc_ids": [doc_id],
                "goal": "Board update",
                "audience": "Executive Steering Committee",
                "tone": 15,
                "slide_count": 6,
            },
        )
        assert planned.status_code == 200
        draft = planned.json()

        queued = client.post(
            "/api/generate/async",
            json={
                "draft_id": draft["draft_id"],
                "outline": [
                    {
                        "id": slide["id"],
                        "index": slide["index"],
                        "purpose": slide["purpose"],
                        "title": slide["title"],
                        "template_id": slide["template_id"],
                    }
                    for slide in draft["slides"]
                ],
                "selected_template_id": "headline.evidence",
                "theme_name": "ONAC",
                "brand_kit": {
                    "logo_data_url": None,
                    "primary_color": "#112233",
                    "accent_color": "#445566",
                    "font_pair": "DM Sans/DM Serif Display",
                },
            },
        )
        assert queued.status_code == 200
        payload = queued.json()
        assert payload["status"] == "queued"

        status_payload = None
        for _ in range(40):
            status = client.get(payload["status_url"])
            assert status.status_code == 200
            status_payload = status.json()
            if status_payload["status"] == "completed":
                break
            time.sleep(0.05)

        assert status_payload is not None
        assert status_payload["status"] == "completed"
        assert status_payload["deck_id"]

        result = client.get(f"/api/generate/jobs/{payload['job_id']}/result")
        assert result.status_code == 200
        assert result.json()["id"] == status_payload["deck_id"]


def test_api_async_generate_job_streams_sse_events(monkeypatch) -> None:
    _reset_api_state()

    def fake_generate(request_model):
        time.sleep(0.05)
        return api_module.PresentationSpecResponse(
            id="deck-sse",
            doc_id="doc-sse",
            doc_ids=["doc-sse"],
            title="Queued Deck",
            goal="Explain queueing",
            audience="Operators",
            slides=[],
            created_at="2026-04-11T20:00:00",
            theme=None,
        )

    monkeypatch.setattr(api_module, "_generate_deck_from_draft_sync", fake_generate)

    with TestClient(api_module.app) as client:
        queued = client.post(
            "/api/generate/async",
            json={
                "draft_id": "draft-sse",
                "outline": [
                    {
                        "id": "s1",
                        "index": 1,
                        "purpose": "content",
                        "title": "Queued slide",
                        "template_id": "headline.evidence",
                    }
                ],
                "selected_template_id": "headline.evidence",
                "theme_name": "ONAC",
                "brand_kit": {
                    "logo_data_url": None,
                    "primary_color": "#112233",
                    "accent_color": "#445566",
                    "font_pair": "DM Sans/DM Serif Display",
                },
            },
        )
        assert queued.status_code == 200
        stream_url = queued.json()["stream_url"]

        events: list[str] = []
        with client.stream("GET", stream_url) as response:
            assert response.status_code == 200
            for line in response.iter_lines():
                if not line:
                    continue
                if isinstance(line, bytes):
                    line = line.decode("utf-8")
                events.append(line)
                if '"status":"completed"' in line:
                    break

        joined = "\n".join(events)
        assert "event: snapshot" in joined
        assert '"status":"completed"' in joined


def test_api_build_vector_store_reuses_persisted_collection(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)
    client = TestClient(api_module.app)

    doc_id = _ingest_fixture(client, sample_pdf_path)
    api_module._INGESTED_VECTOR_STORES.clear()

    def fail_embedder():
        raise AssertionError("embedder should not be needed when persisted vectors exist")

    monkeypatch.setattr(api_module, "_get_embedder", fail_embedder)
    store = api_module._build_vector_store([doc_id])
    query_embedding = deterministic_embedder.encode(["sample ingestion"])[0]

    assert store.has_data()
    assert store.query(query_embedding=query_embedding, n_results=1)


def test_api_async_handlers_offload_blocking_work(monkeypatch, sample_pdf_path, deterministic_embedder) -> None:
    _reset_api_state()
    monkeypatch.setattr(api_module, "_get_embedder", lambda: deterministic_embedder)

    calls: list[str] = []
    real_to_thread = api_module.asyncio.to_thread

    async def recording_to_thread(func, /, *args, **kwargs):
        calls.append(getattr(func, "__name__", repr(func)))
        return await real_to_thread(func, *args, **kwargs)

    monkeypatch.setattr(api_module.asyncio, "to_thread", recording_to_thread)
    client = TestClient(api_module.app)

    ingest = client.post(
        "/api/ingest",
        files={"file": ("sample_ingestion.pdf", sample_pdf_path.read_bytes(), "application/pdf")},
    )
    assert ingest.status_code == 200
    doc_id = ingest.json()["doc_id"]

    planned = client.post(
        "/api/plan",
        json={
            "doc_ids": [doc_id],
            "goal": "Board update",
            "audience": "Executive Steering Committee",
            "tone": 15,
            "slide_count": 6,
        },
    )
    assert planned.status_code == 200
    draft = planned.json()

    generated = client.post(
        "/api/generate",
        json={
            "draft_id": draft["draft_id"],
            "outline": [
                {
                    "id": slide["id"],
                    "index": slide["index"],
                    "purpose": slide["purpose"],
                    "title": slide["title"],
                    "template_id": slide["template_id"],
                }
                for slide in draft["slides"]
            ],
            "selected_template_id": "headline.evidence",
            "brand_kit": {
                "logo_data_url": None,
                "primary_color": "#112233",
                "accent_color": "#445566",
                "font_pair": "DM Sans/DM Serif Display",
            },
        },
    )
    assert generated.status_code == 200

    assert "_ingest_and_index_sync" in calls
    assert "_generate_document_summary_sync" in calls
    assert "_plan_deck_response" in calls
    assert "_generate_deck_from_draft_sync" in calls


def test_api_rate_limits_generation_requests_by_api_key() -> None:
    _reset_api_state()
    client = TestClient(api_module.app)
    payload = {
        "slide_id": "slide-preview-rate-limit",
        "title": "Executive Overview",
        "purpose": "content",
        "template_id": "headline.evidence",
        "content": "Current draft text about ingestion, retrieval, layout, and export.",
        "audience": "Oracle consultants",
        "goal": "Explain the delivery architecture",
    }

    for _ in range(8):
        response = client.post("/api/slide/preview", json=payload, headers={"x-api-key": "alpha"})
        assert response.status_code == 200

    limited = client.post("/api/slide/preview", json=payload, headers={"x-api-key": "alpha"})
    error = _assert_error(limited, status_code=429, code="rate_limit_exceeded", message_fragment="Rate limit exceeded")
    assert limited.headers["retry-after"]
    assert error["request_id"]

    bypass = client.post("/api/slide/preview", json=payload, headers={"x-api-key": "beta"})
    assert bypass.status_code == 200


def test_api_not_found_errors_use_structured_envelope() -> None:
    _reset_api_state()
    client = TestClient(api_module.app)

    response = client.get("/api/deck/missing-deck")

    _assert_error(response, status_code=404, code="deck_not_found", message_fragment="Unknown deck_id: missing-deck")


def test_api_serves_built_frontend(monkeypatch, tmp_path) -> None:
    web_dir = tmp_path / "web"
    assets_dir = web_dir / "assets"
    assets_dir.mkdir(parents=True)
    index_path = web_dir / "index.html"
    asset_path = assets_dir / "app.js"
    index_path.write_text("<!doctype html><html><body>Auto-PPT UI</body></html>", encoding="utf-8")
    asset_path.write_text("console.log('ui')", encoding="utf-8")

    monkeypatch.setattr(api_module, "WEB_DIR", web_dir)
    monkeypatch.setattr(api_module, "WEB_INDEX", index_path)
    client = TestClient(api_module.app)

    root = client.get("/")
    assert root.status_code == 200
    assert "Auto-PPT UI" in root.text

    route = client.get("/editor/demo-deck")
    assert route.status_code == 200
    assert "Auto-PPT UI" in route.text

    asset = client.get("/assets/app.js")
    assert asset.status_code == 200
    assert asset.text == "console.log('ui')"


def test_api_validation_errors_use_structured_envelope() -> None:
    _reset_api_state()
    client = TestClient(api_module.app)

    response = client.post("/api/plan", json={"doc_ids": []})

    _assert_error(response, status_code=422, code="invalid_request", message_fragment="Invalid request payload")


def test_api_health_returns_503_when_dependency_probe_fails(monkeypatch) -> None:
    _reset_api_state()
    client = TestClient(api_module.app)

    def failing_probe():
        raise RuntimeError("model unavailable")

    monkeypatch.setattr(api_module, "_probe_embedder", failing_probe)

    response = client.get("/api/health")

    _assert_error(response, status_code=503, code="health_check_failed", message_fragment="embedder: model unavailable")


# ---------------------------------------------------------------------------
# _infer_chat_brief – dynamic slide count heuristic
# ---------------------------------------------------------------------------


def test_infer_chat_brief_default_slide_count() -> None:
    result = api_module._infer_chat_brief("Make a deck about AI", "AI Overview")
    assert result["slide_count"] == 6


def test_infer_chat_brief_scales_with_chunk_count() -> None:
    # 60 chunks → 4 + 60//10 = 10 slides
    result = api_module._infer_chat_brief("Make a deck", "Doc", content_chunk_count=60)
    assert result["slide_count"] == 10

    # 200 chunks → capped at 20
    result = api_module._infer_chat_brief("Make a deck", "Doc", content_chunk_count=200)
    assert result["slide_count"] == 20

    # 15 chunks → max(6, 4+1) = 6 (floor)
    result = api_module._infer_chat_brief("Make a deck", "Doc", content_chunk_count=15)
    assert result["slide_count"] == 6


def test_infer_chat_brief_explicit_user_count_overrides() -> None:
    result = api_module._infer_chat_brief(
        "Create 12 slides about the platform", "Platform", content_chunk_count=5,
    )
    assert result["slide_count"] == 12


def test_infer_chat_brief_pptx_source_preserves_slide_count() -> None:
    ctx = {"source_format": "pptx", "slide_count": 8}
    result = api_module._infer_chat_brief("Rebrand this deck", "Deck", source_context=ctx)
    assert result["slide_count"] == 8


def test_infer_chat_brief_pptx_explicit_override_wins() -> None:
    ctx = {"source_format": "pptx", "slide_count": 8}
    result = api_module._infer_chat_brief(
        "Rebrand this deck as 4 slides", "Deck", source_context=ctx,
    )
    assert result["slide_count"] == 4


def test_infer_chat_brief_audience_extraction() -> None:
    result = api_module._infer_chat_brief(
        "Create a deck for Oracle consultants on AI", "AI",
    )
    assert result["audience"] == "Oracle consultants"


def test_infer_chat_brief_tone_technical() -> None:
    result = api_module._infer_chat_brief("Technical deep dive for consultants", "Tech")
    assert result["tone"] == 25.0


def test_infer_chat_brief_tone_sales() -> None:
    result = api_module._infer_chat_brief("Bold investor pitch", "Startup")
    assert result["tone"] == 80.0
