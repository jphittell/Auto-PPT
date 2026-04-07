from __future__ import annotations

from collections.abc import Callable
import base64
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from pptx_gen.ingestion.schemas import (
    ContentElementType,
    ContentObject,
    DocumentInfo,
    IngestionOptions,
    IngestionRequest,
    SourceInfo,
    SourceType,
)


class DeterministicEmbedder:
    def encode(self, texts: list[str]) -> list[list[float]]:
        vectors: list[list[float]] = []
        for text in texts:
            encoded = text.encode("utf-8")
            total = sum(encoded)
            length = len(encoded) or 1
            vectors.append(
                [
                    float(length),
                    float(total % 101),
                    float(sum(encoded[::2]) % 103),
                    float(sum(encoded[1::2]) % 107),
                    float(text.count(" ") + 1),
                    float(sum(1 for char in text if char.isdigit())),
                    float(sum(1 for char in text.lower() if char in "aeiou")),
                    float(ord(text[0]) if text else 0),
                ]
            )
        return vectors


@pytest.fixture()
def sample_pdf_path() -> Path:
    return Path(__file__).parent / "fixtures" / "sample_ingestion.pdf"


@pytest.fixture()
def deterministic_embedder() -> DeterministicEmbedder:
    return DeterministicEmbedder()


@pytest.fixture()
def sample_ingestion_request() -> IngestionRequest:
    return IngestionRequest(
        source=SourceInfo(type=SourceType.UPLOAD, id="src-sample", uri="/tmp/sample.txt"),
        document=DocumentInfo(
            title="Sample",
            mime_type="text/plain",
            language="en",
            elements=[
                ContentObject(
                    doc_id="sample-doc",
                    element_id="e0001",
                    page=1,
                    type=ContentElementType.TITLE,
                    text="Sample Deck",
                ),
                ContentObject(
                    doc_id="sample-doc",
                    element_id="e0002",
                    page=1,
                    type=ContentElementType.PARAGRAPH,
                    text="Contact me at sample@example.com for the weekly report.",
                ),
                ContentObject(
                    doc_id="sample-doc",
                    element_id="e0003",
                    page=1,
                    type=ContentElementType.PARAGRAPH,
                    text="Contact me at sample@example.com for the weekly report.",
                ),
            ],
        ),
        options=IngestionOptions(max_chunk_chars=1200),
    )


@pytest.fixture()
def style_tokens_payload() -> dict:
    return {
        "fonts": {
            "heading": "Aptos Display",
            "body": "Aptos",
            "mono": "Cascadia Code",
        },
        "colors": {
            "bg": "#FFFFFF",
            "text": "#111111",
            "accent": "#0A84FF",
            "muted": "#6B7280",
        },
        "spacing": {
            "margin_in": 0.5,
            "gutter_in": 0.25,
        },
        "images": {
            "source_policy": "provided_only",
            "style_prompt": "clean, editorial, brand-consistent",
        },
    }


@pytest.fixture()
def make_citation() -> Callable[[], dict]:
    def factory() -> dict:
        return {
            "source_id": "doc-finance",
            "locator": "doc-finance:page2",
            "quote": "Revenue reached $29.3M.",
            "confidence": 0.93,
        }

    return factory


@pytest.fixture()
def make_block(make_citation: Callable[[], dict]) -> Callable[..., dict]:
    def factory(
        *,
        block_id: str = "b1",
        kind: str = "text",
        content: dict | None = None,
        with_citation: bool = True,
    ) -> dict:
        payload = content or {"text": "Revenue increased six percent quarter over quarter."}
        return {
            "block_id": block_id,
            "kind": kind,
            "content": payload,
            "source_citations": [make_citation()] if with_citation else [],
            "asset_refs": [],
        }

    return factory


@pytest.fixture()
def make_slide(make_block: Callable[..., dict]) -> Callable[..., dict]:
    def factory(
        *,
        slide_id: str = "s1",
        purpose: str = "content",
        blocks: list[dict] | None = None,
        template_key: str = "headline.evidence",
    ) -> dict:
        return {
            "slide_id": slide_id,
            "purpose": purpose,
            "layout_intent": {
                "template_key": template_key,
                "strict_template": True,
            },
            "headline": "Performance highlight",
            "speaker_notes": "",
            "blocks": blocks or [make_block()],
        }

    return factory


@pytest.fixture()
def make_presentation_spec(style_tokens_payload: dict, make_slide: Callable[..., dict]) -> Callable[..., dict]:
    def factory(*, slides: list[dict] | None = None, questions_for_user: list[str] | None = None) -> dict:
        return {
            "title": "Quarterly Business Review",
            "audience": "Executive leadership",
            "language": "en-US",
            "theme": {
                "name": "Executive Blue",
                "style_tokens": style_tokens_payload,
            },
            "slides": slides or [make_slide()],
            "questions_for_user": questions_for_user if questions_for_user is not None else [],
        }

    return factory


@pytest.fixture()
def tiny_png_bytes() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Wn0N0sAAAAASUVORK5CYII="
    )


@pytest.fixture()
def make_docx_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        from docx import Document

        path = tmp_path / "sample.docx"
        document = Document()
        document.add_heading("Sample Document", level=0)
        document.add_heading("Overview", level=1)
        document.add_paragraph("This document explains the new operating model.")
        document.add_paragraph("First bullet", style="List Bullet")
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Option"
        table.cell(0, 1).text = "Fit"
        table.cell(1, 0).text = "A"
        table.cell(1, 1).text = "High"
        document.save(path)
        return path

    return factory


@pytest.fixture()
def make_csv_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        path = tmp_path / "sample.csv"
        path.write_text("Option,Fit\nA,High\nB,Medium\n", encoding="utf-8")
        return path

    return factory


@pytest.fixture()
def make_markdown_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        path = tmp_path / "sample.md"
        path.write_text("# Overview\n\nThis is a concise markdown brief.\n\n- First point\n- Second point\n", encoding="utf-8")
        return path

    return factory


@pytest.fixture()
def make_xlsx_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        from openpyxl import Workbook

        path = tmp_path / "sample.xlsx"
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Decision Matrix"
        sheet.append(["Option", "Cost", "Effort"])
        sheet.append(["A", "Low", "Medium"])
        sheet.append(["B", "Medium", "Low"])
        workbook.save(path)
        workbook.close()
        return path

    return factory


@pytest.fixture()
def make_pptx_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        path = tmp_path / "sample.pptx"
        presentation = Presentation()

        title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        title_slide.shapes.title.text = "Sample Deck"
        title_slide.placeholders[1].text = "Executive overview"

        content_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        title_box = content_slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(6), Inches(0.6))
        title_box.text = "Decision Summary"
        body_box = content_slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(5.5), Inches(2.5))
        body_box.text_frame.text = "Primary recommendation"
        bullet = body_box.text_frame.add_paragraph()
        bullet.text = "First point"
        bullet.level = 1
        table_shape = content_slide.shapes.add_table(2, 2, Inches(6.5), Inches(1.2), Inches(4.5), Inches(1.5))
        table = table_shape.table
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Cost"
        table.cell(1, 1).text = "Low"

        presentation.save(path)
        return path

    return factory


@pytest.fixture()
def make_json_file(tmp_path: Path) -> Callable[[], Path]:
    def factory() -> Path:
        path = tmp_path / "sample.json"
        path.write_text(
            '{\n'
            '  "title": "Chart Spec Catalog",\n'
            '  "owner": "Analytics",\n'
            '  "charts": [\n'
            '    {"name": "Revenue", "chart_type": "bar", "priority": 1},\n'
            '    {"name": "Pipeline", "chart_type": "line", "priority": 2}\n'
            '  ],\n'
            '  "tags": ["finance", "forecast"]\n'
            '}\n',
            encoding="utf-8",
        )
        return path

    return factory


@pytest.fixture()
def make_image_file(tmp_path: Path, tiny_png_bytes: bytes) -> Callable[[str], Path]:
    def factory(extension: str = ".png") -> Path:
        extension = extension.lower()
        path = tmp_path / f"sample{extension}"
        if extension == ".png":
            path.write_bytes(tiny_png_bytes)
            return path

        from PIL import Image

        image = Image.new("RGB", (4, 3), color=(12, 34, 56))
        image.save(path)
        return path

    return factory
