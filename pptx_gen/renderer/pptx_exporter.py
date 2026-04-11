"""Deterministic PPTX exporter built on renderer slide operations."""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

from pptx_gen.layout.templates import get_template_definition
from pptx_gen.layout.schemas import PageSize, ResolvedDeckLayout, ResolvedElement, ResolvedElementKind, StyleTokens
from pptx_gen.renderer.slide_ops import (
    add_bullets,
    add_chart_image,
    add_image,
    add_shape,
    add_table,
    add_text,
    add_title,
    extract_local_asset_path,
    extract_table_content,
    extract_text_lines,
    fill_bullets_placeholder,
    fill_picture_placeholder,
    fill_text_placeholder,
    set_background_color,
)


SLIDE_SIZES = {
    PageSize.WIDESCREEN: (13.333, 7.5),
    PageSize.STANDARD: (10.0, 7.5),
}


def export_pptx(
    *,
    layout: ResolvedDeckLayout,
    style_tokens: StyleTokens,
    output_path: str | Path,
    template_path: str | Path | None = None,
) -> Path:
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(template_path)) if template_path else Presentation()
    if template_path:
        _clear_existing_slides(presentation)
    slide_width, slide_height = SLIDE_SIZES[layout.page_size]
    presentation.slide_width = Inches(slide_width)
    presentation.slide_height = Inches(slide_height)

    for slide_layout in layout.slides:
        template_def = _template_definition_for_slide(slide_layout)
        if template_path and template_def and template_def.layout_index is not None:
            slide = presentation.slides.add_slide(presentation.slide_layouts[template_def.layout_index])
        else:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            set_background_color(slide, style_tokens.colors.bg)
        for element in sorted(slide_layout.elements, key=lambda item: item.z):
            if template_path and template_def is not None and _render_element_into_placeholder(slide, element, template_def, style_tokens):
                continue
            _render_element(slide, element, style_tokens)

    presentation.save(str(output_path))
    return output_path


def _clear_existing_slides(presentation: Presentation) -> None:
    """Remove template/example slides while preserving masters and layouts."""
    while len(presentation.slides) > 0:
        slide_id = presentation.slides._sldIdLst[0]
        presentation.part.drop_rel(slide_id.rId)
        del presentation.slides._sldIdLst[0]


def _template_definition_for_slide(slide_layout: Any):
    for element in slide_layout.elements:
        payload = element.payload or {}
        template_key = payload.get("template_key")
        if isinstance(template_key, str):
            return get_template_definition(template_key)
    return None


def _render_element_into_placeholder(slide: Any, element: ResolvedElement, template_def: Any, style_tokens: StyleTokens) -> bool:
    payload = element.payload or {}
    slot_id = payload.get("slot_id")
    if not isinstance(slot_id, str):
        return False
    slot = next((candidate for candidate in template_def.slots if candidate.slot_id == slot_id), None)
    if slot is None or slot.placeholder_idx is None:
        return False
    placeholder = _placeholder_by_idx(slide, slot.placeholder_idx)
    if placeholder is None:
        return False

    content = payload.get("content")
    try:
        if element.kind is ResolvedElementKind.TEXTBOX:
            lines = extract_text_lines(content)
            if payload.get("block_kind") == "bullets" or (isinstance(content, dict) and isinstance(content.get("items"), list)):
                fill_bullets_placeholder(placeholder, lines, style_tokens=style_tokens, style_ref=element.style_ref)
            else:
                fill_text_placeholder(placeholder, "\n".join(lines), style_tokens=style_tokens, style_ref=element.style_ref)
            return True

        if element.kind is ResolvedElementKind.IMAGE:
            image_path = extract_local_asset_path(content)
            if image_path is None:
                raise FileNotFoundError(f"image element {element.element_id} is missing a local asset path")
            fill_picture_placeholder(placeholder, image_path)
            return True

        if element.kind is ResolvedElementKind.SHAPE:
            text = "\n".join(extract_text_lines(content))
            if text:
                fill_text_placeholder(placeholder, text, style_tokens=style_tokens, style_ref=element.style_ref)
                return True
    except Exception:
        return False
    return False


def _placeholder_by_idx(slide: Any, idx: int) -> Any | None:
    for placeholder in slide.placeholders:
        try:
            if placeholder.placeholder_format.idx == idx:
                return placeholder
        except Exception:
            continue
    return None


def _render_element(slide: Any, element: ResolvedElement, style_tokens: StyleTokens) -> None:
    payload = element.payload or {}
    content = payload.get("content")

    if element.kind is ResolvedElementKind.TEXTBOX:
        lines = extract_text_lines(content)
        if payload.get("block_kind") == "bullets" or (isinstance(content, dict) and isinstance(content.get("items"), list)):
            add_bullets(
                slide,
                lines,
                x=element.x,
                y=element.y,
                w=element.w,
                h=element.h,
                style_tokens=style_tokens,
                style_ref=element.style_ref,
            )
        else:
            text = "\n".join(lines)
            renderer = add_title if element.style_ref == "headline" else add_text
            renderer(
                slide,
                text,
                x=element.x,
                y=element.y,
                w=element.w,
                h=element.h,
                style_tokens=style_tokens,
                style_ref=element.style_ref,
            )
        return

    if element.kind is ResolvedElementKind.IMAGE:
        image_path = extract_local_asset_path(content)
        if image_path is None:
            raise FileNotFoundError(f"image element {element.element_id} is missing a local asset path")
        add_image(slide, image_path, x=element.x, y=element.y, w=element.w, h=element.h)
        return

    if element.kind is ResolvedElementKind.CHART:
        image_path = extract_local_asset_path(content)
        if image_path is None:
            raise FileNotFoundError(f"chart element {element.element_id} is missing a local asset path")
        add_chart_image(slide, image_path, x=element.x, y=element.y, w=element.w, h=element.h)
        return

    if element.kind is ResolvedElementKind.TABLE:
        table_content = extract_table_content(content)
        if table_content is None:
            raise ValueError(f"table element {element.element_id} is missing columns/rows content")
        columns, rows = table_content
        add_table(
            slide,
            columns,
            rows,
            x=element.x,
            y=element.y,
            w=element.w,
            h=element.h,
            style_tokens=style_tokens,
            style_ref=element.style_ref,
        )
        return

    if element.kind is ResolvedElementKind.SHAPE:
        text = "\n".join(extract_text_lines(content)) or None
        add_shape(
            slide,
            x=element.x,
            y=element.y,
            w=element.w,
            h=element.h,
            style_tokens=style_tokens,
            style_ref=element.style_ref,
            text=text,
        )
        return

    raise ValueError(f"unsupported resolved element kind: {element.kind}")
