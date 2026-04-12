"""Deterministic PPTX authoring helpers used by the exporter and QA."""

from __future__ import annotations

from pathlib import Path
from typing import Any
from urllib.parse import urlparse

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from pptx_gen.layout.schemas import ResolvedElementKind, StyleTokens
from pptx_gen.renderer.markdown_strip import strip_markdown


_TEXTBOX_KINDS = {ResolvedElementKind.TEXTBOX, ResolvedElementKind.SHAPE}


def add_title(
    slide: Slide,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    return add_text(
        slide,
        text,
        x=x,
        y=y,
        w=w,
        h=h,
        style_tokens=style_tokens,
        style_ref=style_ref or "headline",
    )


def add_text(
    slide: Slide,
    text: str,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Pt(4)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    profile = style_profile(style_tokens, style_ref or "body", ResolvedElementKind.TEXTBOX)
    run = _ensure_first_run(paragraph)
    run.font.name = profile["font_name"]
    run.font.size = Pt(profile["font_size_pt"])
    run.font.bold = profile["bold"]
    run.font.italic = profile["italic"]
    run.font.color.rgb = rgb_from_hex(profile["text_color"])
    return shape


def fill_text_placeholder(
    shape: Any,
    text: str,
    *,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"placeholder {getattr(shape, 'name', 'unknown')} does not support text")
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    profile = style_profile(style_tokens, style_ref or "body", ResolvedElementKind.TEXTBOX)
    run = _ensure_first_run(paragraph)
    run.font.name = profile["font_name"]
    run.font.size = Pt(profile["font_size_pt"])
    run.font.bold = profile["bold"]
    run.font.italic = profile["italic"]
    run.font.color.rgb = rgb_from_hex(profile["text_color"])
    return shape


def add_bullets(
    slide: Slide,
    items: list[str],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    text_frame.margin_left = Pt(4)
    text_frame.margin_right = Pt(4)
    text_frame.margin_top = Pt(2)
    text_frame.margin_bottom = Pt(2)

    profile = style_profile(style_tokens, style_ref or "body", ResolvedElementKind.TEXTBOX)
    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"• {item}"
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        runs = paragraph.runs or (_ensure_first_run(paragraph),)
        for run in runs:
            run.font.name = profile["font_name"]
            run.font.size = Pt(profile["font_size_pt"])
            run.font.bold = profile["bold"]
            run.font.italic = profile["italic"]
            run.font.color.rgb = rgb_from_hex(profile["text_color"])
    return shape


def fill_bullets_placeholder(
    shape: Any,
    items: list[str],
    *,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"placeholder {getattr(shape, 'name', 'unknown')} does not support text")
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    profile = style_profile(style_tokens, style_ref or "body", ResolvedElementKind.TEXTBOX)
    for index, item in enumerate(items):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"\u2022 {item}"
        paragraph.level = 0
        paragraph.alignment = PP_ALIGN.LEFT
        runs = paragraph.runs or (_ensure_first_run(paragraph),)
        for run in runs:
            run.font.name = profile["font_name"]
            run.font.size = Pt(profile["font_size_pt"])
            run.font.bold = profile["bold"]
            run.font.italic = profile["italic"]
            run.font.color.rgb = rgb_from_hex(profile["text_color"])
    return shape


def add_image(slide: Slide, image_path: str | Path, *, x: float, y: float, w: float, h: float):
    path = ensure_local_asset_path(image_path)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def add_chart_image(slide: Slide, image_path: str | Path, *, x: float, y: float, w: float, h: float):
    return add_image(slide, image_path, x=x, y=y, w=w, h=h)


def fill_picture_placeholder(shape: Any, image_path: str | Path):
    path = ensure_local_asset_path(image_path)
    if not hasattr(shape, "insert_picture"):
        raise ValueError(f"placeholder {getattr(shape, 'name', 'unknown')} does not support pictures")
    return shape.insert_picture(str(path))


def add_table(
    slide: Slide,
    columns: list[str],
    rows: list[list[str]],
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
):
    row_count = max(1, len(rows) + 1)
    col_count = max(1, len(columns))
    shape = slide.shapes.add_table(row_count, col_count, Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    profile = style_profile(style_tokens, style_ref or "table", ResolvedElementKind.TABLE)

    for col_index, column_name in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = str(column_name)
        _style_cell(cell, profile, bold=True)

    for row_index, row in enumerate(rows, start=1):
        for col_index in range(col_count):
            cell = table.cell(row_index, col_index)
            value = row[col_index] if col_index < len(row) else ""
            cell.text = str(value)
            _style_cell(cell, profile, bold=False)

    return shape


def add_shape(
    slide: Slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    style_tokens: StyleTokens,
    style_ref: str | None = None,
    text: str | None = None,
):
    style_ref = style_ref or "shape"
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.RECTANGLE
        if style_ref in {"accent_bar", "takeaway"}
        else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    profile = style_profile(style_tokens, style_ref, ResolvedElementKind.SHAPE)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb_from_hex(profile["fill_color"])
    shape.line.color.rgb = rgb_from_hex(profile["line_color"])
    if text:
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text
        run = _ensure_first_run(paragraph)
        run.font.name = profile["font_name"]
        run.font.size = Pt(profile["font_size_pt"])
        run.font.bold = profile["bold"]
        run.font.color.rgb = rgb_from_hex(profile["text_color"])
    return shape


def set_background_color(slide: Slide, color_hex: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_from_hex(color_hex)


def style_profile(
    style_tokens: StyleTokens,
    style_ref: str,
    kind: ResolvedElementKind,
) -> dict[str, Any]:
    colors = style_tokens.colors
    fonts = style_tokens.fonts

    profile: dict[str, Any] = {
        "font_name": fonts.body,
        "font_size_pt": 18,
        "bold": False,
        "italic": False,
        "text_color": colors.text,
        "fill_color": colors.bg,
        "line_color": colors.muted,
    }

    if style_ref == "headline":
        profile.update({"font_name": fonts.heading, "font_size_pt": 28, "bold": True})
    elif style_ref == "subtitle":
        profile.update({"font_size_pt": 16, "text_color": colors.muted})
    elif style_ref in {"meta", "footer", "citation"}:
        profile.update({"font_size_pt": 10, "text_color": colors.muted})
    elif style_ref in {"body", "agenda", "dense_body"}:
        profile.update({"font_size_pt": 14 if style_ref == "dense_body" else 18})
    elif style_ref == "kpi":
        profile.update({"font_size_pt": 24, "bold": True})
    elif style_ref in {"card", "takeaway"}:
        profile.update(
            {
                "font_size_pt": 16,
                "bold": style_ref == "takeaway",
                "fill_color": colors.accent if style_ref == "takeaway" else "#F8FAFC",
                "line_color": colors.accent,
                "text_color": "#FFFFFF" if style_ref == "takeaway" else "#1E293B",
            }
        )
    elif style_ref == "accent_bar":
        profile.update({"fill_color": colors.accent, "line_color": colors.accent, "text_color": "#FFFFFF"})
    elif kind is ResolvedElementKind.TABLE:
        profile.update({"font_size_pt": 12})

    return profile


def ensure_local_asset_path(image_path: str | Path) -> Path:
    path = Path(image_path)
    parsed = urlparse(str(image_path))
    if parsed.scheme in {"http", "https"}:
        raise ValueError("remote asset URLs are not allowed in renderer helpers")
    if not path.exists():
        raise FileNotFoundError(f"local asset does not exist: {path}")
    return path


def extract_text_lines(content: Any) -> list[str]:
    if content is None:
        return []
    lines: list[str]
    if isinstance(content, str):
        lines = [content]
    elif isinstance(content, list):
        lines = [str(item) for item in content if item is not None]
    elif isinstance(content, dict):
        if isinstance(content.get("items"), list):
            lines = [str(item) for item in content["items"] if item is not None]
        else:
            ordered_keys = ("title", "text", "value", "label", "delta", "sub_label", "subtitle", "presenter", "date")
            values = [str(content[key]) for key in ordered_keys if content.get(key)]
            if values:
                lines = values
            else:
                lines = [str(value) for value in content.values() if value not in (None, "", [], {})]
    else:
        lines = [str(content)]
    return [strip_markdown(line) for line in lines]


def extract_table_content(content: Any) -> tuple[list[str], list[list[str]]] | None:
    if not isinstance(content, dict):
        return None
    columns = content.get("columns")
    rows = content.get("rows")
    if not isinstance(columns, list) or not isinstance(rows, list):
        return None
    normalized_rows = []
    for row in rows:
        if isinstance(row, list):
            normalized_rows.append([strip_markdown(str(cell)) for cell in row])
    return [strip_markdown(str(column)) for column in columns], normalized_rows


def extract_local_asset_path(content: Any) -> Path | None:
    candidates = []
    if isinstance(content, str):
        candidates.append(content)
    elif isinstance(content, dict):
        for key in ("local_path", "path", "file_path", "asset_path", "uri"):
            value = content.get(key)
            if isinstance(value, str):
                candidates.append(value)
    for candidate in candidates:
        parsed = urlparse(candidate)
        if parsed.scheme in {"http", "https"}:
            raise ValueError("remote asset URLs are not allowed in renderer helpers")
        path = Path(candidate)
        if path.exists():
            return path
    return None


def rgb_from_hex(value: str) -> RGBColor:
    return RGBColor.from_string(value.removeprefix("#"))


def _style_cell(cell: Any, profile: dict[str, Any], *, bold: bool) -> None:
    paragraph = cell.text_frame.paragraphs[0]
    if not paragraph.runs:
        paragraph.text = cell.text
    run = _ensure_first_run(paragraph)
    run.font.name = profile["font_name"]
    run.font.size = Pt(profile["font_size_pt"])
    run.font.bold = bold
    run.font.color.rgb = rgb_from_hex(profile["text_color"])


def _ensure_first_run(paragraph: Any) -> Any:
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()
