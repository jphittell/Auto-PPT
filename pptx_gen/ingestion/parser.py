"""Structure-first parsing for local files."""

from __future__ import annotations

import mimetypes
import json
import re
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx_gen.ingestion.schemas import (
    ContentElementType,
    ContentObject,
    DocumentInfo,
    IngestionOptions,
    IngestionRequest,
    SourceInfo,
    SourceType,
)


class PDFDependencyError(RuntimeError):
    """Raised when PDF parsing dependencies are unavailable."""


@dataclass(slots=True)
class _FallbackMetadata:
    page_number: int


@dataclass(slots=True)
class _FallbackElement:
    text: str
    metadata: _FallbackMetadata
    category: str = ""
    extensions: dict[str, Any] | None = None


class Title(_FallbackElement):
    pass


class ListItem(_FallbackElement):
    pass


class NarrativeText(_FallbackElement):
    pass


class Table(_FallbackElement):
    pass


class Header(_FallbackElement):
    pass


class Figure(_FallbackElement):
    pass


@dataclass(slots=True)
class _PartitionResult:
    elements: list[Any]
    extensions: dict[str, Any] | None = None


def parse_source(
    source_path: str | Path,
    *,
    title: str | None = None,
    language: str = "en",
    options: IngestionOptions | None = None,
) -> IngestionRequest:
    path = Path(source_path).expanduser().resolve()
    if not path.exists():
        raise FileNotFoundError(f"source file not found: {path}")

    options = options or IngestionOptions()
    mime_type = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    doc_id = _make_doc_id(path)
    source_id = _make_source_id(path)
    partitioned = _partition_file(path)
    raw_elements = partitioned.elements

    normalized_elements: list[ContentObject] = []
    title_claimed = False
    for index, raw_element in enumerate(raw_elements, start=1):
        text = _extract_text(raw_element)
        if not text:
            continue

        page_number = _extract_page_number(raw_element)
        element_type = _map_element_type(raw_element, page_number, title_claimed)
        if element_type is ContentElementType.TITLE:
            title_claimed = True

        normalized_elements.append(
            ContentObject(
                doc_id=doc_id,
                element_id=f"e{index:04d}",
                page=page_number,
                type=element_type,
                text=text,
                extensions=_extract_extensions(raw_element),
            )
        )

    if not normalized_elements:
        raise ValueError(f"no parseable content found in {path}")

    document_title = title or _infer_title(path, normalized_elements)
    return IngestionRequest(
        source=SourceInfo(type=SourceType.UPLOAD, id=source_id, uri=str(path)),
        document=DocumentInfo(
            title=document_title,
            mime_type=mime_type,
            language=language,
            elements=normalized_elements,
        ),
        options=options,
        extensions=partitioned.extensions,
    )


def _partition_file(path: Path) -> _PartitionResult:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _wrap_partition(_partition_pdf(path))
    if suffix in {".txt", ".md"}:
        return _wrap_partition(_partition_text(path))
    if suffix == ".docx":
        return _wrap_partition(_partition_docx(path))
    if suffix == ".csv":
        return _wrap_partition(_partition_csv(path))
    if suffix == ".json":
        return _wrap_partition(_partition_json(path))
    if suffix in {".xlsx", ".xlsm"}:
        return _wrap_partition(_partition_xlsx(path))
    if suffix == ".pptx":
        return _partition_pptx(path)
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        return _wrap_partition(_partition_image(path))
    raise ValueError(f"unsupported source type: {path.suffix}")


def _wrap_partition(value: list[Any] | _PartitionResult) -> _PartitionResult:
    if isinstance(value, _PartitionResult):
        return value
    return _PartitionResult(elements=value)


def _partition_pdf(path: Path) -> list[Any]:
    try:
        from unstructured.partition.pdf import partition_pdf
    except Exception:
        return _partition_pdf_fallback(path)

    try:
        return partition_pdf(filename=str(path), strategy="fast")
    except Exception:
        return _partition_pdf_fallback(path)


def _partition_text(path: Path) -> list[Any]:
    from unstructured.partition.text import partition_text

    return partition_text(filename=str(path))


def _partition_docx(path: Path) -> list[Any]:
    from docx import Document
    from docx.text.paragraph import Paragraph
    from docx.table import Table as DocxTable

    document = Document(str(path))
    elements: list[Any] = []
    for block in _iter_docx_blocks(document):
        if isinstance(block, Paragraph):
            text = re.sub(r"\s+", " ", block.text or "").strip()
            if not text:
                continue
            style_name = str(getattr(block.style, "name", "")).lower()
            if style_name.startswith("title"):
                elements.append(Title(text=text, metadata=_FallbackMetadata(page_number=1)))
            elif "heading" in style_name:
                elements.append(Header(text=text, metadata=_FallbackMetadata(page_number=1)))
            elif "list" in style_name or re.match(r"^[-*•]\s+", text):
                elements.append(
                    ListItem(
                        text=re.sub(r"^[-*•]\s+", "", text).strip(),
                        metadata=_FallbackMetadata(page_number=1),
                    )
                )
            else:
                elements.append(NarrativeText(text=text, metadata=_FallbackMetadata(page_number=1)))
        elif isinstance(block, DocxTable):
            table_text = _stringify_rows([[cell.text for cell in row.cells] for row in block.rows])
            if table_text:
                elements.append(Table(text=table_text, metadata=_FallbackMetadata(page_number=1)))
    return elements


def _pptx_source_extensions(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slide_blueprint: list[dict[str, Any]] = []
    slide_type_counts: Counter[str] = Counter()

    for slide_index, slide in enumerate(presentation.slides, start=1):
        _title_shape, title_text = _find_pptx_slide_title(slide)
        body_texts, bullet_texts, table_count, picture_count, chart_count = _summarize_pptx_slide(slide, exclude_shape=_title_shape)
        slide_type, purpose_hint, template_hint = _classify_pptx_slide(
            slide_index=slide_index,
            slide_count=len(presentation.slides),
            title_text=title_text,
            body_texts=body_texts,
            bullet_texts=bullet_texts,
            table_count=table_count,
            picture_count=picture_count,
            chart_count=chart_count,
        )
        slide_type_counts[slide_type] += 1
        slide_blueprint.append(
            {
                "slide_number": slide_index,
                "title": title_text or f"Slide {slide_index}",
                "slide_type": slide_type,
                "purpose_hint": purpose_hint,
                "template_hint": template_hint,
                "text_preview": _trim_words(" ".join([*body_texts[:2], *bullet_texts[:3]]), 24),
                "text_count": len(body_texts),
                "bullet_count": len(bullet_texts),
                "table_count": table_count,
                "picture_count": picture_count,
                "chart_count": chart_count,
            }
        )

    return {
        "source_format": "pptx",
        "slide_count": len(presentation.slides),
        "slide_types": dict(sorted(slide_type_counts.items())),
        "slide_blueprint": slide_blueprint,
    }


def _find_pptx_slide_title(slide: Any) -> tuple[Any | None, str]:
    """Return (title_shape, title_text) for a pptx slide.

    Prefers the formal title placeholder.  When none is present, falls back to
    the first textbox whose text looks like a heading (single short line near
    the top of the slide).
    """
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        text = re.sub(r"\s+", " ", title_shape.text or "").strip()
        if text:
            return title_shape, text

    # Fallback: pick the topmost single-line text shape as the title.
    best_shape = None
    best_text = ""
    best_top = float("inf")
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = re.sub(r"\s+", " ", shape.text_frame.text or "").strip()
        if not text or "\n" in text.strip():
            continue
        top = getattr(shape, "top", float("inf"))
        if top < best_top and len(text.split()) <= 12:
            best_top = top
            best_shape = shape
            best_text = text

    return best_shape, best_text


def _summarize_pptx_slide(slide: Any, *, exclude_shape: Any = None) -> tuple[list[str], list[str], int, int, int]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    body_texts: list[str] = []
    bullet_texts: list[str] = []
    table_count = 0
    picture_count = 0
    chart_count = 0

    for shape in slide.shapes:
        if exclude_shape is not None and shape is exclude_shape:
            continue
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = re.sub(r"\s+", " ", paragraph.text or "").strip()
                if not text:
                    continue
                if paragraph.level > 0 or re.match(r"^[-*â€¢]\s+", text):
                    bullet_texts.append(re.sub(r"^[-*â€¢]\s+", "", text).strip())
                else:
                    body_texts.append(text)
            continue
        if getattr(shape, "has_table", False):
            table_count += 1
            continue
        if getattr(shape, "has_chart", False):
            chart_count += 1
            continue
        if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
            picture_count += 1
    return body_texts, bullet_texts, table_count, picture_count, chart_count


def _classify_pptx_slide(
    *,
    slide_index: int,
    slide_count: int,
    title_text: str,
    body_texts: list[str],
    bullet_texts: list[str],
    table_count: int,
    picture_count: int,
    chart_count: int,
) -> tuple[str, str, str]:
    lowered_title = title_text.lower().strip()
    short_points = [text for text in [*bullet_texts, *body_texts] if 2 <= len(text.split()) <= 10]

    if slide_index == 1:
        return "title", "title", "title.cover"
    if any(term in lowered_title for term in ("thank you", "questions", "q&a", "next steps", "action", "actions")):
        return "closing", "closing", "closing.actions"
    if chart_count:
        return "chart", "content", "chart.takeaway"
    if table_count:
        return "table", "content", "compare.2col"
    if any(term in lowered_title for term in ("overview", "summary", "executive")):
        return "summary", "content", "exec.summary"
    if len(short_points) >= 4:
        return "matrix", "content", "compare.2col"
    if len(short_points) == 3:
        return "three_point", "content", "exec.summary"
    if picture_count:
        return "photo", "content", "headline.evidence"
    if slide_index == slide_count:
        return "closing", "closing", "closing.actions"
    return "content", "content", "headline.evidence"


def _trim_words(text: str, max_words: int) -> str:
    words = [word for word in str(text).split() if word]
    return " ".join(words[:max_words])


def _partition_csv(path: Path) -> list[Any]:
    import csv

    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.reader(handle)
        rows = [[cell.strip() for cell in row] for row in reader if any(cell.strip() for cell in row)]

    elements: list[Any] = [Title(text=path.stem.replace("_", " "), metadata=_FallbackMetadata(page_number=1))]
    if rows:
        elements.append(Table(text=_stringify_rows(rows), metadata=_FallbackMetadata(page_number=1)))
    return elements


def _partition_xlsx(path: Path) -> list[Any]:
    from openpyxl import load_workbook

    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    elements: list[Any] = [Title(text=path.stem.replace("_", " "), metadata=_FallbackMetadata(page_number=1))]
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        elements.append(Header(text=sheet.title, metadata=_FallbackMetadata(page_number=sheet_index)))
        rows: list[list[str]] = []
        for row in sheet.iter_rows(values_only=True):
            normalized = ["" if value is None else str(value).strip() for value in row]
            if any(normalized):
                rows.append(normalized)
        if rows:
            elements.append(Table(text=_stringify_rows(rows), metadata=_FallbackMetadata(page_number=sheet_index)))
    workbook.close()
    return elements


def _partition_pptx(path: Path) -> _PartitionResult:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = Presentation(str(path))
    elements: list[Any] = []
    title_claimed = False

    for slide_index, slide in enumerate(presentation.slides, start=1):
        title_shape, title_text = _find_pptx_slide_title(slide)
        if title_text:
            cls = Title if not title_claimed else Header
            elements.append(
                cls(
                    text=title_text,
                    metadata=_FallbackMetadata(page_number=slide_index),
                    extensions={
                        "slide_index": slide_index,
                        "slide_title": title_text,
                        "shape_role": "title",
                    },
                )
            )
            title_claimed = True

        for shape in slide.shapes:
            if title_shape is not None and shape is title_shape:
                continue
            if getattr(shape, "has_text_frame", False):
                paragraphs = []
                for paragraph in shape.text_frame.paragraphs:
                    text = re.sub(r"\s+", " ", paragraph.text or "").strip()
                    if text:
                        paragraphs.append((paragraph.level, text))
                for level, text in paragraphs:
                    if level > 0 or re.match(r"^[-*•]\s+", text):
                        elements.append(
                            ListItem(
                                text=re.sub(r"^[-*•]\s+", "", text).strip(),
                                metadata=_FallbackMetadata(page_number=slide_index),
                            )
                        )
                    else:
                        elements.append(NarrativeText(text=text, metadata=_FallbackMetadata(page_number=slide_index)))
            elif getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                table_text = _stringify_rows(rows)
                if table_text:
                    elements.append(Table(text=table_text, metadata=_FallbackMetadata(page_number=slide_index)))
            elif getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                elements.append(Figure(text=shape.name or "Figure", metadata=_FallbackMetadata(page_number=slide_index)))

    return _PartitionResult(elements=elements, extensions=_pptx_source_extensions(path))


def _partition_image(path: Path) -> list[Any]:
    from PIL import Image

    stem_title = path.stem.replace("_", " ").replace("-", " ").strip() or path.name
    with Image.open(path) as image:
        width, height = image.size
        image_format = (image.format or path.suffix.removeprefix(".")).upper()
        mode = image.mode

    return [
        Title(text=stem_title, metadata=_FallbackMetadata(page_number=1)),
        Figure(
            text=f"Image asset {stem_title} ({image_format}, {width}x{height}, mode {mode})",
            metadata=_FallbackMetadata(page_number=1),
            extensions={
                "path": str(path),
                "format": image_format,
                "width_px": width,
                "height_px": height,
                "mode": mode,
            },
        ),
    ]


def _partition_json(path: Path) -> list[Any]:
    with path.open("r", encoding="utf-8-sig") as handle:
        payload = json.load(handle)

    title = _json_document_title(path, payload)
    elements: list[Any] = [Title(text=title, metadata=_FallbackMetadata(page_number=1))]
    _append_json_elements(elements, payload, page_number=1, heading_level=0, current_key=None)
    return elements


def _extract_text(raw_element: Any) -> str:
    value = getattr(raw_element, "text", None)
    if value is None:
        value = str(raw_element)
    normalized = re.sub(r"\s+", " ", value or "").strip()
    return normalized


def _extract_extensions(raw_element: Any) -> dict[str, Any] | None:
    value = getattr(raw_element, "extensions", None)
    return value if isinstance(value, dict) else None


def _partition_pdf_fallback(path: Path) -> list[Any]:
    try:
        from pypdf import PdfReader
    except Exception as exc:  # pragma: no cover - environment specific
        raise PDFDependencyError(
            "PDF parsing requires either the Unstructured PDF stack or the lightweight pypdf fallback."
        ) from exc

    reader = PdfReader(str(path))
    elements: list[Any] = []
    title_claimed = False

    for page_number, page in enumerate(reader.pages, start=1):
        page_text = re.sub(r"\r\n?", "\n", page.extract_text() or "")
        lines = [line.strip() for line in page_text.splitlines() if line.strip()]
        if not lines:
            continue

        paragraph_buffer: list[str] = []

        def flush_paragraph() -> None:
            if not paragraph_buffer:
                return
            elements.append(
                NarrativeText(
                    text=" ".join(paragraph_buffer).strip(),
                    metadata=_FallbackMetadata(page_number=page_number),
                )
            )
            paragraph_buffer.clear()

        for line in lines:
            if not title_claimed:
                elements.append(
                    Title(
                        text=line,
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
                title_claimed = True
                continue

            if re.match(r"^[-*•]\s+", line):
                flush_paragraph()
                elements.append(
                    ListItem(
                        text=re.sub(r"^[-*•]\s+", "", line).strip(),
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
                continue

            if len(line.split()) <= 10 and line == line.title():
                flush_paragraph()
                elements.append(
                    Title(
                        text=line,
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
                continue

            paragraph_buffer.append(line)

        flush_paragraph()

    if not elements:
        raise PDFDependencyError("PDF fallback could not extract text from the file.")
    return elements


def _extract_page_number(raw_element: Any) -> int | None:
    metadata = getattr(raw_element, "metadata", None)
    page_number = getattr(metadata, "page_number", None)
    if page_number is None:
        return None
    try:
        return int(page_number)
    except (TypeError, ValueError):
        return None


def _map_element_type(raw_element: Any, page_number: int | None, title_claimed: bool) -> ContentElementType:
    name = type(raw_element).__name__.lower()
    category = str(getattr(raw_element, "category", "")).lower()
    signature = f"{name} {category}"

    if "figurecaption" in signature or "caption" in signature:
        return ContentElementType.CAPTION
    if "image" in signature or "figure" in signature:
        return ContentElementType.FIGURE
    if "table" in signature:
        return ContentElementType.TABLE
    if "listitem" in signature or "list_item" in signature:
        return ContentElementType.LIST_ITEM
    if "title" in signature:
        if not title_claimed and (page_number in {None, 1}):
            return ContentElementType.TITLE
        return ContentElementType.HEADING
    if "header" in signature or "heading" in signature:
        return ContentElementType.HEADING
    return ContentElementType.PARAGRAPH


def _infer_title(path: Path, elements: list[ContentObject]) -> str:
    stem = path.stem.replace("_", " ").replace("-", " ").strip()
    for element in elements:
        if element.type is ContentElementType.TITLE:
            text = element.text.strip()
            # Prefer the filename stem when the parsed title looks like a fragment
            # (ends with a dangling conjunction/preposition, or is shorter than stem)
            last_word = text.rsplit(" ", 1)[-1].lower().rstrip(".,;:")
            is_fragment = last_word in {"and", "or", "but", "the", "a", "an", "of", "in", "on", "at", "by", "for", "with", "to"}
            if stem and (is_fragment or (len(stem) > len(text) + 5)):
                return stem
            return text
    return stem or path.name


def _make_doc_id(path: Path) -> str:
    slug = _slugify(path.stem)
    return slug or "document"


def _make_source_id(path: Path) -> str:
    slug = _slugify(path.stem)
    return f"src-{slug or 'document'}"


def _slugify(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "-", value.strip()).strip("-_.").lower()


def _stringify_rows(rows: list[list[str]]) -> str:
    cleaned_rows = []
    for row in rows:
        cleaned = [re.sub(r"\s+", " ", cell or "").strip() for cell in row]
        if any(cleaned):
            cleaned_rows.append(" | ".join(cleaned))
    return "\n".join(cleaned_rows).strip()


def _iter_docx_blocks(document: Any):
    from docx.document import Document as DocumentObject
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph

    parent = document.element.body if isinstance(document, DocumentObject) else document
    for child in parent.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, document)
        elif isinstance(child, CT_Tbl):
            yield DocxTable(child, document)


def _append_json_elements(
    elements: list[Any],
    value: Any,
    *,
    page_number: int,
    heading_level: int,
    current_key: str | None,
) -> None:
    if isinstance(value, dict):
        items = list(value.items())
        if current_key and heading_level >= 0:
            elements.append(Header(text=_humanize_key(current_key), metadata=_FallbackMetadata(page_number=page_number)))
        if _looks_like_table_rows(value):
            rows = [[key, _scalar_to_text(item)] for key, item in items]
            elements.append(
                Table(
                    text=_stringify_rows([["Field", "Value"], *rows]),
                    metadata=_FallbackMetadata(page_number=page_number),
                    extensions={"json_shape": "object_table"},
                )
            )
            return
        for key, item in items:
            if _is_scalar(item):
                elements.append(
                    NarrativeText(
                        text=f"{_humanize_key(key)}: {_scalar_to_text(item)}",
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
            else:
                _append_json_elements(
                    elements,
                    item,
                    page_number=page_number,
                    heading_level=heading_level + 1,
                    current_key=key,
                )
        return

    if isinstance(value, list):
        if current_key:
            elements.append(Header(text=_humanize_key(current_key), metadata=_FallbackMetadata(page_number=page_number)))
        if _is_tabular_list(value):
            rows = _tabular_rows_from_list(value)
            if rows:
                elements.append(
                    Table(
                        text=_stringify_rows(rows),
                        metadata=_FallbackMetadata(page_number=page_number),
                        extensions={"json_shape": "list_table"},
                    )
                )
            return
        if all(_is_scalar(item) for item in value):
            for item in value:
                elements.append(
                    ListItem(
                        text=_scalar_to_text(item),
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
            return
        for index, item in enumerate(value, start=1):
            label = f"{current_key or 'Item'} {index}"
            if _is_scalar(item):
                elements.append(
                    NarrativeText(
                        text=f"{label}: {_scalar_to_text(item)}",
                        metadata=_FallbackMetadata(page_number=page_number),
                    )
                )
            else:
                _append_json_elements(
                    elements,
                    item,
                    page_number=page_number,
                    heading_level=heading_level + 1,
                    current_key=label,
                )
        return

    scalar_text = _scalar_to_text(value)
    if current_key:
        elements.append(
            NarrativeText(
                text=f"{_humanize_key(current_key)}: {scalar_text}",
                metadata=_FallbackMetadata(page_number=page_number),
            )
        )
    else:
        elements.append(NarrativeText(text=scalar_text, metadata=_FallbackMetadata(page_number=page_number)))


def _json_document_title(path: Path, payload: Any) -> str:
    if isinstance(payload, dict):
        for key in ("title", "name", "document_title"):
            value = payload.get(key)
            if isinstance(value, str) and value.strip():
                return value.strip()
    return path.stem.replace("_", " ").replace("-", " ").strip() or path.name


def _is_scalar(value: Any) -> bool:
    return value is None or isinstance(value, (str, int, float, bool))


def _scalar_to_text(value: Any) -> str:
    if value is None:
        return "null"
    if isinstance(value, bool):
        return "true" if value else "false"
    return str(value)


def _humanize_key(key: str) -> str:
    cleaned = re.sub(r"[_-]+", " ", str(key)).strip()
    return cleaned[:1].upper() + cleaned[1:] if cleaned else "Field"


def _is_tabular_list(values: list[Any]) -> bool:
    dict_items = [item for item in values if isinstance(item, dict)]
    if len(dict_items) != len(values) or not dict_items:
        return False
    key_sets = [tuple(item.keys()) for item in dict_items]
    first = key_sets[0]
    return all(keys == first for keys in key_sets) and all(_is_scalar(item.get(key)) for item in dict_items for key in first)


def _tabular_rows_from_list(values: list[dict[str, Any]]) -> list[list[str]]:
    if not values:
        return []
    columns = [str(key) for key in values[0].keys()]
    rows = [columns]
    for item in values:
        rows.append([_scalar_to_text(item.get(column)) for column in columns])
    return rows


def _looks_like_table_rows(value: dict[str, Any]) -> bool:
    return bool(value) and all(_is_scalar(item) for item in value.values())
